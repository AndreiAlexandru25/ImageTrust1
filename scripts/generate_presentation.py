#!/usr/bin/env python3
"""
Generate academic PowerPoint presentations for ImageTrust thesis defense.
Produces two files: Romanian (RO) and English (EN), ~38 slides each.

Usage:
    python scripts/generate_presentation.py

Output:
    outputs/ImageTrust_Presentation_RO.pptx
    outputs/ImageTrust_Presentation_EN.pptx
"""

import os
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
FIGURES_DIR = ROOT / "outputs" / "phase3" / "figures"
OUTPUT_DIR = ROOT / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Theme colors
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ===========================================================================
# Helper functions
# ===========================================================================

def set_slide_bg(slide, color=WHITE):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, height=Inches(1.2)):
    """Add a dark blue rectangle at the top as a title bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    return shape


def add_footer_bar(slide, slide_num, total_slides):
    """Add a thin footer bar with slide number."""
    bar_h = Inches(0.35)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_HEIGHT - bar_h, SLIDE_WIDTH, bar_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total_slides}"
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_slide(slide, title_text, bullets, slide_num, total,
                     start_top=Inches(1.5)):
    """Standard content slide with title bar and bullet points."""
    set_slide_bg(slide, WHITE)
    add_title_bar(slide)
    add_footer_bar(slide, slide_num, total)

    # Title text in the bar
    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8),
        title_text, font_size=28, bold=True, color=WHITE, font_name="Calibri"
    )

    # Bullet content
    txBox = slide.shapes.add_textbox(
        Inches(0.8), start_top, Inches(11.5), Inches(5.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        # Support sub-bullets with "  - " prefix
        if bullet.startswith("  - "):
            p.level = 1
            text = bullet[4:]
        else:
            p.level = 0
            text = bullet.lstrip("- ")

        run = p.add_run()
        run.text = text
        run.font.size = Pt(16) if p.level == 0 else Pt(14)
        run.font.color.rgb = DARK_GRAY if p.level == 0 else MEDIUM_GRAY
        run.font.name = "Calibri"

    return txBox


def add_table_slide(slide, title_text, headers, rows, slide_num, total,
                    col_widths=None):
    """Slide with a PowerPoint table."""
    set_slide_bg(slide, WHITE)
    add_title_bar(slide)
    add_footer_bar(slide, slide_num, total)

    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8),
        title_text, font_size=28, bold=True, color=WHITE, font_name="Calibri"
    )

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(11.5) if col_widths is None else sum(col_widths)
    tbl_left = Inches(0.9)
    tbl_top = Inches(1.6)
    row_height = Inches(0.4)
    tbl_height = row_height * n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
    )
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Calibri"
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            bg = LIGHT_GRAY if i % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = DARK_GRAY
                    run.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_figure_slide(slide, title_text, fig_path, slide_num, total,
                     caption=None, img_top=Inches(1.5), img_height=Inches(5.0)):
    """Slide with title bar and centered image."""
    set_slide_bg(slide, WHITE)
    add_title_bar(slide)
    add_footer_bar(slide, slide_num, total)

    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8),
        title_text, font_size=28, bold=True, color=WHITE, font_name="Calibri"
    )

    if fig_path and os.path.exists(fig_path):
        # Center the image
        from PIL import Image
        img = Image.open(fig_path)
        w, h = img.size
        aspect = w / h
        img_w = img_height * aspect
        max_w = Inches(11.5)
        if img_w > max_w:
            img_w = max_w
            img_height = img_w / aspect
        img_left = (SLIDE_WIDTH - img_w) / 2
        slide.shapes.add_picture(str(fig_path), img_left, img_top, img_w, img_height)
    else:
        add_text_box(
            slide, Inches(2), Inches(3), Inches(9), Inches(1),
            f"[Figure: {fig_path}]", font_size=14, color=MEDIUM_GRAY,
            alignment=PP_ALIGN.CENTER
        )

    if caption:
        add_text_box(
            slide, Inches(1), SLIDE_HEIGHT - Inches(0.9), Inches(11), Inches(0.4),
            caption, font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER
        )


def add_two_column_slide(slide, title_text, left_bullets, right_bullets,
                         slide_num, total, left_title=None, right_title=None):
    """Slide with two columns of bullet points."""
    set_slide_bg(slide, WHITE)
    add_title_bar(slide)
    add_footer_bar(slide, slide_num, total)

    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8),
        title_text, font_size=28, bold=True, color=WHITE, font_name="Calibri"
    )

    col_top = Inches(1.5)
    col_w = Inches(5.5)

    for col_idx, (bullets, col_title) in enumerate([
        (left_bullets, left_title), (right_bullets, right_title)
    ]):
        left_pos = Inches(0.6) if col_idx == 0 else Inches(7.0)

        if col_title:
            add_text_box(
                slide, left_pos, col_top, col_w, Inches(0.4),
                col_title, font_size=18, bold=True, color=ACCENT_BLUE
            )
            bullet_top = col_top + Inches(0.5)
        else:
            bullet_top = col_top

        txBox = slide.shapes.add_textbox(left_pos, bullet_top, col_w, Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            if bullet.startswith("  - "):
                p.level = 1
                text = bullet[4:]
            else:
                p.level = 0
                text = bullet.lstrip("- ")
            run = p.add_run()
            run.text = text
            run.font.size = Pt(14) if p.level == 0 else Pt(12)
            run.font.color.rgb = DARK_GRAY if p.level == 0 else MEDIUM_GRAY
            run.font.name = "Calibri"


def add_section_slide(slide, section_title, section_subtitle, slide_num, total):
    """Section divider slide with centered text on dark background."""
    set_slide_bg(slide, DARK_BLUE)
    add_footer_bar(slide, slide_num, total)

    add_text_box(
        slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
        section_title, font_size=40, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )
    if section_subtitle:
        add_text_box(
            slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
            section_subtitle, font_size=20, color=LIGHT_BLUE,
            alignment=PP_ALIGN.CENTER, font_name="Calibri"
        )


# ===========================================================================
# Slide content definitions (bilingual)
# ===========================================================================

def get_slides(lang="EN"):
    """Return list of slide definitions. lang='EN' or 'RO'."""

    # --- Convenience translations ---
    t = {}  # translation dict
    if lang == "RO":
        t = {
            # Part 1
            "title": "ImageTrust: Detectarea Forensica a Imaginilor\nGenerate de Inteligenta Artificiala",
            "subtitle": "Disertatie de Master",
            "author": "Andrei Alexandru",
            "university": "Universitatea Politehnica Timisoara\nFacultatea de Automatica si Calculatoare",
            "date": "Februarie 2026",
            "advisor": "Coordonator: Prof. Dr.",

            "s2_title": "Problema de Cercetare",
            "s2_bullets": [
                "Generatoarele AI (Midjourney, DALL-E, Stable Diffusion) produc imagini foto-realiste",
                "Detectia vizuala devine imposibila chiar si pentru experti",
                "Trei deficiente ale solutiilor existente:",
                "  - 1. Lipsa calibrarii - scorurile nu reflecta probabilitati reale",
                "  - 2. Fara regiune de incertitudine - verdictele sunt binare (real/fals)",
                "  - 3. Robustetea la degradari lipseste - compresie, recadraj, blur",
                "Necesitatea unui sistem forensic complet, calibrat si cu grad de incredere",
            ],

            "s3_title": "Obiective de Cercetare",
            "s3_bullets": [
                "Contributia 1: Fuziune multi-backbone cu meta-clasificare",
                "  - ResNet-50 + EfficientNet-B0 + ViT-B/16 -> vector 4097d -> XGBoost/MLP",
                "Contributia 2: Calibrare post-hoc + predictie conformala",
                "  - Temperature scaling, Platt, Isotonic -> ECE < 1%",
                "  - LAC conformal -> 95.2% acoperire cu 9.3% abstinere",
                "Contributia 3: Evaluare completa pentru publicare",
                "  - Comparatie cu 4 baseline-uri, degradari, cross-generator, ablatie",
                "  - Teste de semnificatie statistica (McNemar, DeLong)",
            ],

            "s4_title": "Scop si Amploare",
            "s4_bullets": [
                "604,376 imagini de antrenament (151K x 4 variante de compresie)",
                "141,000 imagini GenImage (24 generatoare) pentru evaluare cross-generator",
                "6 metode comparate cu acelasi protocol, metrici si split-uri",
                "Aplicatie desktop Windows (.exe) cu PySide6",
                "  - Interfata drag-and-drop, rapoarte PDF, functionare offline",
                "Articol LNCS-ready cu 10 figuri, 7 tabele, 28 referinte",
            ],

            "s5_title": "Cerinte Academice (Mandatorii)",
            "s5_bullets": [
                "Comparatie cu baseline-uri: LogReg, CNN single, ViT, XGBoost multi-backbone",
                "Calibrare: diagrame de fiabilitate, metrici ECE, scalare temperatura",
                "Regiune INCERT/abstinere: predictie conformala LAC/APS/RAPS",
                "Studiu ablativ: contributia fiecarei componente",
                "Model de amenintari + limitari: documentate explicit",
                "Reproductibilitate: seed-uri fixe, split-uri, hiperparametri, spec. hardware",
                "Metrici de eficienta: ms/imagine, imagini/sec, VRAM",
            ],

            # Part 2 - Dataset
            "s6_title": "Set de Date - Prezentare Generala",
            "s6_bullets": [
                "5 surse de date combinate:",
                "  - CIFAKE: 120,000 imagini (60K reale + 60K Stable Diffusion)",
                "  - COCO 2017: 118,287 fotografii reale",
                "  - FFHQ: 70,000 fete reale (Flickr)",
                "  - SFHQ: 89,914 fete sintetice (StyleGAN)",
                "  - Deepfake Faces: 140,000 imagini (70K reale + 70K deepfake)",
                "Total brut: ~538,000 imagini unice",
                "57.9% reale / 42.1% generate AI (echilibrat)",
            ],

            "s7_title": "Variante de Compresie",
            "s7_bullets": [
                "Fiecare imagine procesata in 4 variante:",
                "  - Original: calitate nealterata",
                "  - WhatsApp: JPEG Q=75, resize 1600px, strip EXIF",
                "  - Instagram: JPEG Q=85, resize 1080px, strip EXIF",
                "  - Screenshot: JPEG Q=90, padding 2px, gamma shift",
                "Total: 151,094 x 4 = 604,376 imagini de antrenament",
                "Scopul: robustetea la scenarii reale de distributie",
                "Augmentare implicita fara cost suplimentar",
            ],

            "s8_title": "Partitionarea Datelor",
            "s8_bullets": [
                "Split stratificat: 70% antrenare / 15% validare / 15% test",
                "  - Antrenare: 422,710 imagini",
                "  - Validare: 90,692 imagini",
                "  - Test: 90,974 imagini",
                "Seed fix: 42 (reproductibilitate garantata)",
                "Stratificare pe sursa + eticheta (real/AI)",
                "Nicio scurgere de date intre split-uri (verificat)",
            ],

            "s9_title": "GenImage - Evaluare Cross-Generator",
            "s9_bullets": [
                "Dataset separat: GenImage (141,000 imagini)",
                "24 generatoare AI complet noi (nevazute in antrenament):",
                "  - GAN-uri: BigGAN, StyleGAN, ProGAN, CycleGAN, StarGAN, GauGAN",
                "  - Diffusion: ADM, DDPM, IDDPM, LDM, PNDM, Stable Diffusion v1.4/v1.5",
                "  - Autoregresive: DALL-E, Midjourney, Glide, VQDM",
                "Evaluare zero-shot: modelul NU a vazut aceste generatoare",
                "Test de generalizare: cat de bine extrapoleaza la surse noi?",
            ],

            # Part 3 - Architecture
            "s10_title": "Arhitectura Sistemului",

            "s11_title": "Faza 1: Extractia Embedding-urilor",
            "s11_bullets": [
                "3 backbone-uri CNN/ViT pre-antrenate (ImageNet):",
                "  - ResNet-50: 2048 dimensiuni (features spatiale profunde)",
                "  - EfficientNet-B0: 1280 dimensiuni (eficienta parametrica)",
                "  - ViT-B/16: 768 dimensiuni (atentie globala, patch-uri)",
                "Semnal suplimentar: NIQE (calitate imagine fara referinta)",
                "Vector concatenat: 2048 + 1280 + 768 + 1 = 4097 dimensiuni",
                "Extractie pe GPU: ~14 ms/imagine (RTX 5080)",
                "Embedding-uri salvate pe disc pentru reutilizare",
            ],

            "s12_title": "Faza 2: Meta-Clasificatoare",
            "s12_bullets": [
                "XGBoost Meta-Classifier:",
                "  - Input: 4097 features -> 1000 arbori, adancime 8",
                "  - Acuratete: 88.7%, AUC: 96.0%, ECE: 2.1%",
                "  - Timp inferenta: 0.44 ms/imagine (doar clasificare)",
                "MLP Meta-Classifier:",
                "  - 4097 -> 1024 -> 512 -> 256 -> 1 (ReLU + Dropout 0.3)",
                "  - Acuratete: 89.1%, AUC: 96.3%, ECE: 3.9%",
                "  - Performanta cea mai buna, dar ECE mai mare",
                "Fuziunea > oricare backbone individual",
            ],

            "s13_title": "Calibrare Post-Hoc",
            "s13_bullets": [
                "Scorurile brute NU sunt probabilitati reale",
                "3 metode de calibrare evaluate:",
                "  - Temperature Scaling: un singur parametru T, simplu si eficient",
                "  - Platt Scaling: regresie logistica pe logits",
                "  - Isotonic Regression: non-parametrica, cea mai flexibila",
                "Rezultate calibrare (pe setul de test):",
                "  - Necalibrat: ECE = 2.1% - 3.9%",
                "  - Temperature: ECE = 1.6% (cel mai bun raport simplitate/performanta)",
                "  - Isotonic: ECE = 0.6% - 0.8% (cel mai precis)",
                "ECE masoara distanta medie intre incredere si acuratete reala",
            ],

            "s14_title": "Predictie Conformala",
            "s14_bullets": [
                "Problema: cand modelul NU stie, ar trebui sa spuna 'Nu stiu'",
                "Solutie: predictie conformala (Vovk et al., 2005)",
                "3 metode implementate:",
                "  - LAC (Least Ambiguous set-valued Classifier)",
                "  - APS (Adaptive Prediction Sets)",
                "  - RAPS (Regularized Adaptive Prediction Sets)",
                "Rezultate LAC la alpha=0.05:",
                "  - Acoperire: 95.2% (garantie teoretica: 95%)",
                "  - Abstinere: 9.3% imagini marcate INCERT",
                "  - Prag: 0.7652 (calculat pe 604K esantioane)",
                "Verdict: REAL / AI-GENERAT / INCERT",
            ],

            "s15_title": "Tehnici de Antrenament",
            "s15_bullets": [
                "AMP (Automatic Mixed Precision): antrenament FP16 -> 2x viteza",
                "SWA (Stochastic Weight Averaging): mediere ponderi ultimele epoci",
                "Mixup (alpha=0.4): interpolare intre esantioane",
                "Label Smoothing (0.1): regularizare soft a etichetelor",
                "RandAugment: augmentari automate (rotatie, crop, color jitter)",
                "Multi-seed training: 3 seed-uri x 3 backbone-uri = 9 modele",
                "EMA (Exponential Moving Average): mediere exponentiala a ponderilor",
                "Cosine warmup scheduler: crestere graduala lr -> descrestere cosinus",
            ],

            "s16_title": "Aplicatia Desktop",
            "s16_bullets": [
                "Interfata PySide6 (Qt6) - aspect modern, profesional",
                "Functionalitati principale:",
                "  - Drag & drop imagini (single sau batch)",
                "  - Analiza completa: detectie + metadata + explicabilitate",
                "  - Verdict cu incredere calibrata + zona INCERT",
                "  - Tab-uri: Rezultate, Heatmap, Metadata, Provenienta",
                "  - Export raport PDF + JSON",
                "Functionare 100% offline (modele locale)",
                "Distribuibil: .exe Windows via PyInstaller",
            ],

            # Part 4 - Results
            "s17_title": "Comparatie Principala (6 Metode)",
            "s18_title": "Curbe ROC - Toate Metodele",
            "s19_title": "Robustetea la Degradari",
            "s20_title": "Evaluare Cross-Generator",
            "s21_title": "Heatmap Cross-Generator",
            "s22_title": "Studiu Ablativ",
            "s23_title": "Rezultate Calibrare",
            "s24_title": "Diagrame de Fiabilitate",
            "s25_title": "Predictie Conformala - Rezultate",
            "s26_title": "Semnificatie Statistica",
            "s27_title": "Metrici de Eficienta",
            "s28_title": "Diagnostic Overfitting",

            "s19_bullets": [
                "4 tipuri de degradari evaluate:",
                "  - Compresie JPEG: Q=50, 70, 85, 95",
                "  - Redimensionare: 25%, 50%, 75%",
                "  - Blur Gaussian: sigma=0.5, 1.0, 2.0",
                "  - Zgomot: 1%, 3%, 5%",
                "Scadere maxima AUC: 0.3% (de la 96.0% la 95.7%)",
                "Robustetea vine din antrenamentul pe 4 variante de compresie",
                "Cel mai robust: XGBoost meta-clasificator (fuziune 3 backbone-uri)",
            ],

            "s20_bullets": [
                "Performanta pe 24 generatoare noi (nevazute):",
                "  - Cele mai bune: Deepfake >99%, COCO >99.9% (surse similare)",
                "  - Mediu: CIFAKE-SD 85.2%, BigGAN 78%, ProGAN 72%",
                "  - Dificile: StarGAN 69.3%, CycleGAN 48%",
                "  - Foarte dificile: DDPM <1%, IDDPM <1% (difuzie avansata)",
                "AUC mediu cross-gen: ~80% (variatie mare intre generatoare)",
                "Limitare onesta: generatoarele difuzionale sunt cele mai dificile",
            ],

            "s25_bullets": [
                "Predictie conformala LAC la alpha=0.05:",
                "  - Acoperire empirica: 95.2% (target: 95.0%)",
                "  - Rata de abstinere: 9.3%",
                "  - Dimensiune medie set: 1.09",
                "La alpha=0.10: acoperire 91.4%, abstinere 0%",
                "APS si RAPS la alpha=0.05: acoperire 100%, dar abstinere 100%",
                "LAC ofera cel mai bun compromis acoperire/precizie",
                "Pragul conformal: 0.7652 (calculat pe 604K esantioane)",
            ],

            "s26_bullets": [
                "McNemar test (perechi de erori):",
                "  - MLP vs LogReg: chi2=536.4, p<0.001 (semnificativ)",
                "  - MLP vs ResNet-50 single: chi2=58.5, p<0.001",
                "  - MLP vs ViT single: chi2=149.2, p<0.001",
                "  - XGBoost vs MLP: chi2=20.8, p<0.001",
                "DeLong test (AUC-uri):",
                "  - Diferentele AUC nu sunt semnificative statistic",
                "  - Clasificarile difera, dar discriminarea e similara",
                "Corectie Bonferroni: alpha=0.01 (5 comparatii)",
            ],

            "s27_bullets": [
                "Pipeline complet (embedding + clasificare):",
                "  - Extractie embedding: 14.0 ms/imagine (3 backbone-uri)",
                "  - XGBoost inferenta: 0.44 ms/imagine",
                "  - MLP inferenta: 0.34 ms/imagine",
                "  - Total: 14.3 ms/imagine (70 imagini/sec)",
                "Hardware: NVIDIA RTX 5080, 16GB VRAM",
                "  - RAM: 32 GB, CPU: 8 cores / 16 threads",
                "Dimensiune model: XGBoost 8.7 MB, MLP 18.6 MB",
                "Suficient pentru analiza in timp real si batch",
            ],

            # Part 5 - Engineering
            "s29_title": "Statistici Codebase",
            "s29_bullets": [
                "~100 fisiere Python sursa",
                "15 module: detection, evaluation, forensics, metadata, etc.",
                "9 fisiere de test (6 unit + 3 integration)",
                "20+ documente: arhitectura, model amenintari, ghid utilizator",
                "~15,000 linii de cod (fara teste si scripturi)",
                "Acoperire teste: detection, calibrare, metrici, metadata",
                "CI/CD: pre-commit hooks, black, isort, flake8, mypy",
            ],

            "s30_title": "Stiva Tehnologica",
            "s31_title": "Pipeline ML",
            "s32_title": "Artefacte de Publicare",

            "s30_bullets_left": [
                "ML & Detectie:",
                "  - PyTorch 2.x + torchvision",
                "  - XGBoost, scikit-learn",
                "  - HuggingFace Transformers",
                "  - ONNX (export optional)",
                "Frontend:",
                "  - PySide6 (Qt6) - desktop",
                "  - Streamlit - web demo",
            ],
            "s30_bullets_right": [
                "Backend:",
                "  - FastAPI + Uvicorn",
                "  - Pydantic v2 (validare)",
                "Infrastructure:",
                "  - Docker + docker-compose",
                "  - PyInstaller (.exe)",
                "  - pytest + coverage",
                "  - WandB (logging optional)",
            ],

            "s31_bullets": [
                "Faza 1: Extractie Embedding-uri",
                "  - 3 backbone-uri -> embedding-uri 4097d -> salvare pe disc",
                "  - Script: scripts/orchestrator/phase1_embeddings.py",
                "Faza 2: Antrenament Meta-Clasificator",
                "  - XGBoost + MLP pe embedding-uri -> calibrare -> conformal",
                "  - Script: scripts/orchestrator/phase2_train.py",
                "Faza 3: Pipeline de Publicare",
                "  - Baseline-uri, degradari, cross-gen, ablatie, figuri, tabele",
                "  - Script: scripts/orchestrator/phase3_publication.py",
                "Totul automatizat: o singura comanda per faza",
            ],

            "s32_bullets": [
                "10 figuri generate automat (matplotlib/seaborn):",
                "  - ROC, calibrare, heatmap, degradari, ablatie, conformal, etc.",
                "7 tabele LaTeX (LNCS format):",
                "  - Comparatie, cross-gen, degradari, ablatie, calibrare, eficienta, conformal",
                "28 referinte bibliografice cu URL-uri clickabile",
                "Articol in format LNCS (Springer) - gata de submisie",
                "Reproducibilitate completa: seed-uri, split-uri, hiperparametri",
            ],

            # Part 6 - Future
            "s33_title": "Model de Amenintari",
            "s33_bullets": [
                "Atacuri adversariale:",
                "  - Perturbatii imperceptibile pot pacali detectorul",
                "  - Evaluare planificata cu FGSM, PGD, C&W",
                "Drift temporal:",
                "  - Generatoarele noi (DALL-E 4, Sora) pot fi diferite",
                "  - Necesita re-antrenare periodica",
                "Bias de domeniu:",
                "  - Antrenat pe fete + scene -> poate esua pe alte domenii",
                "  - SFHQ (StyleGAN) provoaca fals pozitive pe fete reale FFHQ",
                "Limitari de acoperire:",
                "  - Nu detecteaza: inpainting partial, super-rezolutie, editare text",
            ],

            "s34_title": "Evaluare Onesta a Limitarilor",
            "s34_bullets": [
                "Cross-generator: 17-20% din generatoarele noi sunt nedetectate",
                "  - DDPM, IDDPM: <1% acuratete (difuzie avansata)",
                "  - CycleGAN: ~48% (aproape aleator)",
                "SFHQ false pozitive: fetele sintetice seamana cu cele reale",
                "Calibrare: Platt scaling degradeaza uneori (ECE creste)",
                "APS/RAPS la alpha=0.05: abstinere 100% (neutilizabil)",
                "Niciun test adversarial inca (planificat)",
                "Dataset limitat la fete + scene (nu documente, text, medical)",
            ],

            "s35_title": "Lucrari Viitoare - Termen Scurt",
            "s35_bullets": [
                "Grad-CAM si explicabilitate vizuala",
                "  - Heatmap-uri pe regiunile suspecte",
                "  - Analiza patch-uri si frecvente",
                "Validare C2PA (Content Authenticity)",
                "  - Verificare provenienta criptografica",
                "Evaluare adversariala completa",
                "  - FGSM, PGD, patch-uri adversariale",
                "Detectie screenshot / recaptura",
                "  - Clasificator dedicat bazat pe artefacte specifice",
            ],

            "s36_title": "Lucrari Viitoare - Termen Lung",
            "s36_bullets": [
                "Invatare continua (Continual Learning)",
                "  - Adaptare la generatoare noi fara re-antrenare completa",
                "Multi-task learning",
                "  - Detectie + identificare generator + tip manipulare",
                "Federated Learning",
                "  - Antrenament distribuit fara partajarea datelor",
                "Analiza video si audio",
                "  - Extindere de la imagini statice la multimedia",
                "API cloud cu rate limiting",
                "  - Serviciu de detectie ca microserviciu",
            ],

            "s37_title": "Extensii Planificate",
            "s37_bullets": [
                "Detectie screenshot / recaptura",
                "  - Artefacte: padding, gamma shift, UI overlay, moiré",
                "Clasificator retele sociale",
                "  - WhatsApp vs Instagram vs Telegram vs neprocesat",
                "  - Bazat pe semnatura JPEG + metadata strip pattern",
                "Integrare modele HuggingFace (online, optional)",
                "  - 4 modele publice ca al doilea nivel de detectie",
                "  - dima806, umm-maybe, Nahrawy, NYUAD-ComNets",
                "Benchmarking continuu cu generatoare 2026",
            ],

            # Part 7
            "s38_title": "Concluzii",
            "s38_bullets": [
                "ImageTrust: sistem complet de detectie forensica AI",
                "Contributii principale:",
                "  - Fuziune multi-backbone (3 modele) cu meta-clasificare",
                "  - Calibrare post-hoc: ECE redus de la 3.9% la 0.6%",
                "  - Predictie conformala: 95.2% acoperire, 9.3% abstinere",
                "Rezultate: AUC 96.3%, acuratete 89.1%, robustetea <0.3% pierdere",
                "Evaluare completa: baseline-uri, cross-gen, degradari, ablatie",
                "Aplicatie desktop functionala (.exe Windows)",
                "",
                "Intrebari?",
            ],
        }
    else:  # English
        t = {
            "title": "ImageTrust: Forensic Detection of\nAI-Generated Images",
            "subtitle": "Master's Thesis",
            "author": "Andrei Alexandru",
            "university": "Politehnica University Timisoara\nFaculty of Automation and Computing",
            "date": "February 2026",
            "advisor": "Advisor: Prof. Dr.",

            "s2_title": "Research Problem",
            "s2_bullets": [
                "AI generators (Midjourney, DALL-E, Stable Diffusion) produce photo-realistic images",
                "Visual detection becomes impossible even for human experts",
                "Three shortcomings of existing solutions:",
                "  - 1. Lack of calibration - scores don't reflect true probabilities",
                "  - 2. No uncertainty region - verdicts are binary (real/fake)",
                "  - 3. Missing degradation robustness - compression, resize, blur",
                "Need for a complete forensic system with calibration and confidence",
            ],

            "s3_title": "Research Objectives",
            "s3_bullets": [
                "Contribution 1: Multi-backbone fusion with meta-classification",
                "  - ResNet-50 + EfficientNet-B0 + ViT-B/16 -> 4097d vector -> XGBoost/MLP",
                "Contribution 2: Post-hoc calibration + conformal prediction",
                "  - Temperature scaling, Platt, Isotonic -> ECE < 1%",
                "  - LAC conformal -> 95.2% coverage with 9.3% abstention",
                "Contribution 3: Publication-grade evaluation",
                "  - Comparison with 4 baselines, degradation, cross-generator, ablation",
                "  - Statistical significance tests (McNemar, DeLong)",
            ],

            "s4_title": "Project Scope",
            "s4_bullets": [
                "604,376 training images (151K x 4 compression variants)",
                "141,000 GenImage images (24 generators) for cross-generator evaluation",
                "6 methods compared with same protocol, metrics, and splits",
                "Windows desktop application (.exe) with PySide6",
                "  - Drag-and-drop interface, PDF reports, offline operation",
                "LNCS-ready paper with 10 figures, 7 tables, 28 references",
            ],

            "s5_title": "Academic Requirements (Mandatory)",
            "s5_bullets": [
                "Baseline comparison: LogReg, CNN single, ViT, multi-backbone XGBoost",
                "Calibration: reliability diagrams, ECE metrics, temperature scaling",
                "UNCERTAIN/abstain region: LAC/APS/RAPS conformal prediction",
                "Ablation study: per-component contribution analysis",
                "Threat model + limitations: explicitly documented",
                "Reproducibility: fixed seeds, splits, hyperparameters, hardware specs",
                "Efficiency metrics: ms/image, images/sec, VRAM usage",
            ],

            # Part 2
            "s6_title": "Dataset Overview",
            "s6_bullets": [
                "5 combined data sources:",
                "  - CIFAKE: 120,000 images (60K real + 60K Stable Diffusion)",
                "  - COCO 2017: 118,287 real photographs",
                "  - FFHQ: 70,000 real faces (Flickr)",
                "  - SFHQ: 89,914 synthetic faces (StyleGAN)",
                "  - Deepfake Faces: 140,000 images (70K real + 70K deepfake)",
                "Total raw: ~538,000 unique images",
                "57.9% real / 42.1% AI-generated (balanced)",
            ],

            "s7_title": "Compression Variants",
            "s7_bullets": [
                "Each image processed into 4 variants:",
                "  - Original: unaltered quality",
                "  - WhatsApp: JPEG Q=75, resize 1600px, strip EXIF",
                "  - Instagram: JPEG Q=85, resize 1080px, strip EXIF",
                "  - Screenshot: JPEG Q=90, padding 2px, gamma shift",
                "Total: 151,094 x 4 = 604,376 training images",
                "Purpose: robustness to real-world distribution scenarios",
                "Implicit augmentation at zero additional cost",
            ],

            "s8_title": "Data Partitioning",
            "s8_bullets": [
                "Stratified split: 70% train / 15% validation / 15% test",
                "  - Training: 422,710 images",
                "  - Validation: 90,692 images",
                "  - Test: 90,974 images",
                "Fixed seed: 42 (guaranteed reproducibility)",
                "Stratification by source + label (real/AI)",
                "No data leakage between splits (verified)",
            ],

            "s9_title": "GenImage - Cross-Generator Evaluation",
            "s9_bullets": [
                "Separate dataset: GenImage (141,000 images)",
                "24 completely new AI generators (unseen during training):",
                "  - GANs: BigGAN, StyleGAN, ProGAN, CycleGAN, StarGAN, GauGAN",
                "  - Diffusion: ADM, DDPM, IDDPM, LDM, PNDM, Stable Diffusion v1.4/v1.5",
                "  - Autoregressive: DALL-E, Midjourney, Glide, VQDM",
                "Zero-shot evaluation: model has NOT seen these generators",
                "Generalization test: how well does it extrapolate to new sources?",
            ],

            # Part 3
            "s10_title": "System Architecture",

            "s11_title": "Phase 1: Embedding Extraction",
            "s11_bullets": [
                "3 pre-trained CNN/ViT backbones (ImageNet):",
                "  - ResNet-50: 2048 dimensions (deep spatial features)",
                "  - EfficientNet-B0: 1280 dimensions (parameter efficiency)",
                "  - ViT-B/16: 768 dimensions (global attention, patches)",
                "Additional signal: NIQE (no-reference image quality)",
                "Concatenated vector: 2048 + 1280 + 768 + 1 = 4097 dimensions",
                "GPU extraction: ~14 ms/image (RTX 5080)",
                "Embeddings saved to disk for reuse",
            ],

            "s12_title": "Phase 2: Meta-Classifiers",
            "s12_bullets": [
                "XGBoost Meta-Classifier:",
                "  - Input: 4097 features -> 1000 trees, depth 8",
                "  - Accuracy: 88.7%, AUC: 96.0%, ECE: 2.1%",
                "  - Inference time: 0.44 ms/image (classification only)",
                "MLP Meta-Classifier:",
                "  - 4097 -> 1024 -> 512 -> 256 -> 1 (ReLU + Dropout 0.3)",
                "  - Accuracy: 89.1%, AUC: 96.3%, ECE: 3.9%",
                "  - Best overall performance, but higher ECE",
                "Fusion > any single backbone",
            ],

            "s13_title": "Post-Hoc Calibration",
            "s13_bullets": [
                "Raw scores are NOT true probabilities",
                "3 calibration methods evaluated:",
                "  - Temperature Scaling: single parameter T, simple and effective",
                "  - Platt Scaling: logistic regression on logits",
                "  - Isotonic Regression: non-parametric, most flexible",
                "Calibration results (on test set):",
                "  - Uncalibrated: ECE = 2.1% - 3.9%",
                "  - Temperature: ECE = 1.6% (best simplicity/performance ratio)",
                "  - Isotonic: ECE = 0.6% - 0.8% (most precise)",
                "ECE measures average gap between confidence and actual accuracy",
            ],

            "s14_title": "Conformal Prediction",
            "s14_bullets": [
                "Problem: when the model doesn't know, it should say 'I don't know'",
                "Solution: conformal prediction (Vovk et al., 2005)",
                "3 methods implemented:",
                "  - LAC (Least Ambiguous set-valued Classifier)",
                "  - APS (Adaptive Prediction Sets)",
                "  - RAPS (Regularized Adaptive Prediction Sets)",
                "LAC results at alpha=0.05:",
                "  - Coverage: 95.2% (theoretical guarantee: 95%)",
                "  - Abstention: 9.3% of images marked UNCERTAIN",
                "  - Threshold: 0.7652 (computed on 604K samples)",
                "Verdict: REAL / AI-GENERATED / UNCERTAIN",
            ],

            "s15_title": "Training Techniques",
            "s15_bullets": [
                "AMP (Automatic Mixed Precision): FP16 training -> 2x speedup",
                "SWA (Stochastic Weight Averaging): weight averaging over last epochs",
                "Mixup (alpha=0.4): interpolation between samples",
                "Label Smoothing (0.1): soft label regularization",
                "RandAugment: automatic augmentations (rotation, crop, color jitter)",
                "Multi-seed training: 3 seeds x 3 backbones = 9 models",
                "EMA (Exponential Moving Average): exponential weight averaging",
                "Cosine warmup scheduler: gradual lr warmup -> cosine decay",
            ],

            "s16_title": "Desktop Application",
            "s16_bullets": [
                "PySide6 (Qt6) interface - modern, professional look",
                "Main features:",
                "  - Drag & drop images (single or batch)",
                "  - Full analysis: detection + metadata + explainability",
                "  - Verdict with calibrated confidence + UNCERTAIN zone",
                "  - Tabs: Results, Heatmap, Metadata, Provenance",
                "  - PDF + JSON report export",
                "100% offline operation (local models)",
                "Distributable: Windows .exe via PyInstaller",
            ],

            # Part 4
            "s17_title": "Main Comparison (6 Methods)",
            "s18_title": "ROC Curves - All Methods",
            "s19_title": "Degradation Robustness",
            "s20_title": "Cross-Generator Evaluation",
            "s21_title": "Cross-Generator Heatmap",
            "s22_title": "Ablation Study",
            "s23_title": "Calibration Results",
            "s24_title": "Reliability Diagrams",
            "s25_title": "Conformal Prediction Results",
            "s26_title": "Statistical Significance",
            "s27_title": "Efficiency Metrics",
            "s28_title": "Overfitting Diagnostic",

            "s19_bullets": [
                "4 types of degradation evaluated:",
                "  - JPEG compression: Q=50, 70, 85, 95",
                "  - Resize: 25%, 50%, 75%",
                "  - Gaussian blur: sigma=0.5, 1.0, 2.0",
                "  - Noise: 1%, 3%, 5%",
                "Maximum AUC drop: 0.3% (from 96.0% to 95.7%)",
                "Robustness comes from training on 4 compression variants",
                "Most robust: XGBoost meta-classifier (3-backbone fusion)",
            ],

            "s20_bullets": [
                "Performance on 24 new generators (unseen):",
                "  - Best: Deepfake >99%, COCO >99.9% (similar sources)",
                "  - Medium: CIFAKE-SD 85.2%, BigGAN 78%, ProGAN 72%",
                "  - Difficult: StarGAN 69.3%, CycleGAN 48%",
                "  - Very difficult: DDPM <1%, IDDPM <1% (advanced diffusion)",
                "Average cross-gen AUC: ~80% (high variance across generators)",
                "Honest limitation: diffusion generators are most challenging",
            ],

            "s25_bullets": [
                "Conformal prediction LAC at alpha=0.05:",
                "  - Empirical coverage: 95.2% (target: 95.0%)",
                "  - Abstention rate: 9.3%",
                "  - Average set size: 1.09",
                "At alpha=0.10: coverage 91.4%, abstention 0%",
                "APS and RAPS at alpha=0.05: coverage 100%, but abstention 100%",
                "LAC provides the best coverage/precision trade-off",
                "Conformal threshold: 0.7652 (computed on 604K samples)",
            ],

            "s26_bullets": [
                "McNemar test (paired errors):",
                "  - MLP vs LogReg: chi2=536.4, p<0.001 (significant)",
                "  - MLP vs ResNet-50 single: chi2=58.5, p<0.001",
                "  - MLP vs ViT single: chi2=149.2, p<0.001",
                "  - XGBoost vs MLP: chi2=20.8, p<0.001",
                "DeLong test (AUCs):",
                "  - AUC differences are NOT statistically significant",
                "  - Classifications differ, but discrimination is similar",
                "Bonferroni correction: alpha=0.01 (5 comparisons)",
            ],

            "s27_bullets": [
                "Full pipeline (embedding + classification):",
                "  - Embedding extraction: 14.0 ms/image (3 backbones)",
                "  - XGBoost inference: 0.44 ms/image",
                "  - MLP inference: 0.34 ms/image",
                "  - Total: 14.3 ms/image (70 images/sec)",
                "Hardware: NVIDIA RTX 5080, 16GB VRAM",
                "  - RAM: 32 GB, CPU: 8 cores / 16 threads",
                "Model size: XGBoost 8.7 MB, MLP 18.6 MB",
                "Sufficient for real-time and batch analysis",
            ],

            # Part 5
            "s29_title": "Codebase Statistics",
            "s29_bullets": [
                "~100 Python source files",
                "15 modules: detection, evaluation, forensics, metadata, etc.",
                "9 test files (6 unit + 3 integration)",
                "20+ documents: architecture, threat model, user guide",
                "~15,000 lines of code (excluding tests and scripts)",
                "Test coverage: detection, calibration, metrics, metadata",
                "CI/CD: pre-commit hooks, black, isort, flake8, mypy",
            ],

            "s30_title": "Technology Stack",
            "s31_title": "ML Pipeline",
            "s32_title": "Publication Artifacts",

            "s30_bullets_left": [
                "ML & Detection:",
                "  - PyTorch 2.x + torchvision",
                "  - XGBoost, scikit-learn",
                "  - HuggingFace Transformers",
                "  - ONNX (optional export)",
                "Frontend:",
                "  - PySide6 (Qt6) - desktop",
                "  - Streamlit - web demo",
            ],
            "s30_bullets_right": [
                "Backend:",
                "  - FastAPI + Uvicorn",
                "  - Pydantic v2 (validation)",
                "Infrastructure:",
                "  - Docker + docker-compose",
                "  - PyInstaller (.exe)",
                "  - pytest + coverage",
                "  - WandB (optional logging)",
            ],

            "s31_bullets": [
                "Phase 1: Embedding Extraction",
                "  - 3 backbones -> 4097d embeddings -> save to disk",
                "  - Script: scripts/orchestrator/phase1_embeddings.py",
                "Phase 2: Meta-Classifier Training",
                "  - XGBoost + MLP on embeddings -> calibration -> conformal",
                "  - Script: scripts/orchestrator/phase2_train.py",
                "Phase 3: Publication Pipeline",
                "  - Baselines, degradation, cross-gen, ablation, figures, tables",
                "  - Script: scripts/orchestrator/phase3_publication.py",
                "Fully automated: one command per phase",
            ],

            "s32_bullets": [
                "10 auto-generated figures (matplotlib/seaborn):",
                "  - ROC, calibration, heatmap, degradation, ablation, conformal, etc.",
                "7 LaTeX tables (LNCS format):",
                "  - Comparison, cross-gen, degradation, ablation, calibration, efficiency, conformal",
                "28 bibliography references with clickable URLs",
                "Paper in LNCS (Springer) format - submission-ready",
                "Full reproducibility: seeds, splits, hyperparameters",
            ],

            # Part 6
            "s33_title": "Threat Model",
            "s33_bullets": [
                "Adversarial attacks:",
                "  - Imperceptible perturbations can fool the detector",
                "  - Planned evaluation with FGSM, PGD, C&W",
                "Temporal drift:",
                "  - New generators (DALL-E 4, Sora) may differ significantly",
                "  - Requires periodic retraining",
                "Domain bias:",
                "  - Trained on faces + scenes -> may fail on other domains",
                "  - SFHQ (StyleGAN) causes false positives on real FFHQ faces",
                "Coverage limitations:",
                "  - Does not detect: partial inpainting, super-resolution, text editing",
            ],

            "s34_title": "Honest Assessment of Limitations",
            "s34_bullets": [
                "Cross-generator: 17-20% of new generators are undetected",
                "  - DDPM, IDDPM: <1% accuracy (advanced diffusion)",
                "  - CycleGAN: ~48% (nearly random)",
                "SFHQ false positives: synthetic faces resemble real ones",
                "Calibration: Platt scaling sometimes degrades (ECE increases)",
                "APS/RAPS at alpha=0.05: 100% abstention (unusable)",
                "No adversarial testing yet (planned)",
                "Dataset limited to faces + scenes (not documents, text, medical)",
            ],

            "s35_title": "Short-Term Future Work",
            "s35_bullets": [
                "Grad-CAM and visual explainability",
                "  - Heatmaps on suspicious regions",
                "  - Patch and frequency analysis",
                "C2PA validation (Content Authenticity)",
                "  - Cryptographic provenance verification",
                "Complete adversarial evaluation",
                "  - FGSM, PGD, adversarial patches",
                "Screenshot / recapture detection",
                "  - Dedicated classifier based on specific artifacts",
            ],

            "s36_title": "Long-Term Future Work",
            "s36_bullets": [
                "Continual Learning",
                "  - Adapt to new generators without full retraining",
                "Multi-task Learning",
                "  - Detection + generator identification + manipulation type",
                "Federated Learning",
                "  - Distributed training without sharing data",
                "Video and audio analysis",
                "  - Extend from static images to multimedia",
                "Cloud API with rate limiting",
                "  - Detection service as a microservice",
            ],

            "s37_title": "Planned Extensions",
            "s37_bullets": [
                "Screenshot / recapture detection",
                "  - Artifacts: padding, gamma shift, UI overlay, moire",
                "Social media classifier",
                "  - WhatsApp vs Instagram vs Telegram vs unprocessed",
                "  - Based on JPEG signature + metadata strip patterns",
                "HuggingFace model integration (online, optional)",
                "  - 4 public models as second detection tier",
                "  - dima806, umm-maybe, Nahrawy, NYUAD-ComNets",
                "Continuous benchmarking with 2026 generators",
            ],

            # Part 7
            "s38_title": "Conclusions",
            "s38_bullets": [
                "ImageTrust: complete forensic AI image detection system",
                "Main contributions:",
                "  - Multi-backbone fusion (3 models) with meta-classification",
                "  - Post-hoc calibration: ECE reduced from 3.9% to 0.6%",
                "  - Conformal prediction: 95.2% coverage, 9.3% abstention",
                "Results: AUC 96.3%, accuracy 89.1%, robustness <0.3% loss",
                "Complete evaluation: baselines, cross-gen, degradation, ablation",
                "Functional desktop application (Windows .exe)",
                "",
                "Questions?",
            ],
        }

    return t


# ===========================================================================
# Main presentation builder
# ===========================================================================

def build_presentation(lang="EN"):
    """Build the full presentation for the given language."""
    t = get_slides(lang)
    total_slides = 38

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]  # Blank

    # -----------------------------------------------------------------------
    # Slide 1: Title
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_text_box(
        slide, Inches(1), Inches(1.2), Inches(11), Inches(2.0),
        t["title"], font_size=36, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )
    add_text_box(
        slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
        t["subtitle"], font_size=24, color=LIGHT_BLUE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )

    # Horizontal line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5.3), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_text_box(
        slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
        t["author"], font_size=22, color=WHITE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )
    add_text_box(
        slide, Inches(1), Inches(5.2), Inches(11), Inches(0.8),
        t["university"], font_size=16, color=LIGHT_BLUE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )
    add_text_box(
        slide, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
        t["date"], font_size=16, color=LIGHT_BLUE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )

    # -----------------------------------------------------------------------
    # Slide 2: Problem Statement
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s2_title"], t["s2_bullets"], 2, total_slides)

    # -----------------------------------------------------------------------
    # Slide 3: Research Objectives
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s3_title"], t["s3_bullets"], 3, total_slides)

    # -----------------------------------------------------------------------
    # Slide 4: Project Scope
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s4_title"], t["s4_bullets"], 4, total_slides)

    # -----------------------------------------------------------------------
    # Slide 5: Academic Requirements
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s5_title"], t["s5_bullets"], 5, total_slides)

    # -----------------------------------------------------------------------
    # Slide 6: Dataset Overview (section divider)
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    section_sub = "Partea 2: Setul de Date" if lang == "RO" else "Part 2: Dataset"
    add_section_slide(slide, section_sub, "604,376 images | 5 sources | 4 variants",
                      6, total_slides)

    # -----------------------------------------------------------------------
    # Slide 7: Dataset Overview (content - renumbered from plan's slide 6)
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s6_title"], t["s6_bullets"], 7, total_slides)

    # -----------------------------------------------------------------------
    # Slide 8: Compression Variants
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s7_title"], t["s7_bullets"], 8, total_slides)

    # -----------------------------------------------------------------------
    # Slide 9: Data Split
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s8_title"], t["s8_bullets"], 9, total_slides)

    # -----------------------------------------------------------------------
    # Slide 10: GenImage Cross-Gen
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s9_title"], t["s9_bullets"], 10, total_slides)

    # -----------------------------------------------------------------------
    # Slide 11: Architecture section divider
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    section_sub = "Partea 3: Arhitectura" if lang == "RO" else "Part 3: Architecture"
    add_section_slide(slide, section_sub,
                      "3-phase pipeline | Multi-backbone fusion | Conformal prediction",
                      11, total_slides)

    # -----------------------------------------------------------------------
    # Slide 12: System Architecture diagram placeholder
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide)
    add_footer_bar(slide, 12, total_slides)
    add_text_box(
        slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8),
        t["s10_title"], font_size=28, bold=True, color=WHITE
    )

    # Architecture diagram as text boxes
    phases = [
        ("Phase 1\nEmbedding Extraction" if lang == "EN" else "Faza 1\nExtractie Embedding",
         "ResNet-50 (2048d)\nEfficientNet-B0 (1280d)\nViT-B/16 (768d)\n+ NIQE (1d)",
         Inches(0.5)),
        ("Phase 2\nMeta-Classification" if lang == "EN" else "Faza 2\nMeta-Clasificare",
         "4097d vector\nXGBoost (1000 trees)\nMLP (4 layers)\nCalibration",
         Inches(4.7)),
        ("Phase 3\nDecision" if lang == "EN" else "Faza 3\nDecizie",
         "Conformal Prediction\nLAC / APS / RAPS\nREAL | AI | UNCERTAIN",
         Inches(8.9)),
    ]

    for phase_title, phase_desc, left_pos in phases:
        # Phase box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8),
            Inches(3.8), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = phase_title
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

        # Description box
        desc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(3.1),
            Inches(3.8), Inches(2.2)
        )
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = LIGHT_BLUE
        desc_box.line.color.rgb = ACCENT_BLUE
        tf2 = desc_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = phase_desc
        run2.font.size = Pt(13)
        run2.font.color.rgb = DARK_BLUE
        run2.font.name = "Calibri"

    # Arrows between phases
    for arrow_left in [Inches(4.3), Inches(8.5)]:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(2.0),
            Inches(0.5), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE
        arrow.line.fill.background()

    # -----------------------------------------------------------------------
    # Slides 13-17: Architecture details
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s11_title"], t["s11_bullets"], 13, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s12_title"], t["s12_bullets"], 14, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s13_title"], t["s13_bullets"], 15, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s14_title"], t["s14_bullets"], 16, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s15_title"], t["s15_bullets"], 17, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s16_title"], t["s16_bullets"], 18, total_slides)

    # -----------------------------------------------------------------------
    # Slide 19: Results section divider
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    section_sub = "Partea 4: Rezultate" if lang == "RO" else "Part 4: Results"
    add_section_slide(slide, section_sub,
                      "AUC 96.3% | ECE 0.6% | 14.3 ms/image",
                      19, total_slides)

    # -----------------------------------------------------------------------
    # Slide 20: Main comparison table
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    if lang == "RO":
        headers = ["Metoda", "Acc (%)", "F1 (%)", "AUC (%)", "ECE (%)", "Tip"]
    else:
        headers = ["Method", "Acc (%)", "F1 (%)", "AUC (%)", "ECE (%)", "Type"]

    rows = [
        ["B1: LogReg (ResNet-50)", "87.0", "83.9", "93.9", "2.0",
         "Classical" if lang == "EN" else "Clasic"],
        ["B2: XGB (ResNet-50)", "88.4", "85.3", "95.6", "2.2",
         "CNN" if lang == "EN" else "CNN"],
        ["B2: XGB (EffNet-B0)", "88.1", "84.8", "95.3", "2.4",
         "CNN" if lang == "EN" else "CNN"],
        ["B3: XGB (ViT-B/16)", "88.0", "84.6", "95.0", "2.7",
         "ViT" if lang == "EN" else "ViT"],
        ["Ours: XGB (3-backbone)", "88.7", "85.9", "96.0", "2.1",
         "Fusion" if lang == "EN" else "Fuziune"],
        ["Ours: MLP (3-backbone)", "89.1", "86.8", "96.3", "3.9",
         "Fusion" if lang == "EN" else "Fuziune"],
    ]

    col_widths = [Inches(3.5), Inches(1.3), Inches(1.3), Inches(1.3),
                  Inches(1.3), Inches(1.5)]
    add_table_slide(slide, t["s17_title"], headers, rows, 20, total_slides,
                    col_widths=col_widths)

    # -----------------------------------------------------------------------
    # Slide 21: ROC curves figure
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    fig_path = FIGURES_DIR / "fig1_roc_curves.png"
    add_figure_slide(slide, t["s18_title"], fig_path, 21, total_slides,
                     caption="ROC curves for all 6 methods on test set (90,692 samples)")

    # -----------------------------------------------------------------------
    # Slide 22: Degradation robustness
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s19_title"], t["s19_bullets"], 22, total_slides)

    # -----------------------------------------------------------------------
    # Slide 23: Cross-generator table
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    if lang == "RO":
        cg_headers = ["Sursa", "Esantioane", "Acc (%)", "Observatii"]
    else:
        cg_headers = ["Source", "Samples", "Acc (%)", "Notes"]
    cg_rows = [
        ["COCO (real)", "15,656", "99.9",
         "Best" if lang == "EN" else "Cel mai bun"],
        ["Deepfake", "17,920", "99.7",
         "Excellent" if lang == "EN" else "Excelent"],
        ["FFHQ (real)", "4,452", "95.5",
         "Good" if lang == "EN" else "Bun"],
        ["CIFAKE-Real", "4,708", "93.4",
         "Good" if lang == "EN" else "Bun"],
        ["CIFAKE-SD", "3,212", "85.2",
         "Moderate" if lang == "EN" else "Moderat"],
        ["Other (mixed)", "28,092", "66.6",
         "Challenging" if lang == "EN" else "Dificil"],
    ]
    cg_widths = [Inches(3.0), Inches(2.0), Inches(2.0), Inches(3.0)]
    add_table_slide(slide, t["s20_title"], cg_headers, cg_rows, 23, total_slides,
                    col_widths=cg_widths)

    # -----------------------------------------------------------------------
    # Slide 24: Cross-gen heatmap
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    fig_path = FIGURES_DIR / "fig3_cross_source_heatmap.png"
    add_figure_slide(slide, t["s21_title"], fig_path, 24, total_slides)

    # -----------------------------------------------------------------------
    # Slide 25: Ablation study table
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    if lang == "RO":
        abl_headers = ["Configuratie", "Acc (%)", "F1 (%)", "AUC (%)",
                       "ECE (%)", "dF1"]
    else:
        abl_headers = ["Configuration", "Acc (%)", "F1 (%)", "AUC (%)",
                       "ECE (%)", "dF1"]
    abl_rows = [
        ["Full system", "90.8", "88.7", "96.9", "0.012", "---"],
        ["- Quality features", "90.9", "88.8", "96.9", "0.013", "+0.0"],
        ["ResNet-50 only", "90.5", "88.3", "96.4", "0.011", "-0.4"],
        ["EffNet-B0 only", "90.1", "87.8", "96.1", "0.011", "-0.9"],
        ["ViT-B/16 only", "89.4", "86.9", "95.6", "0.011", "-1.8"],
        ["Originals only" if lang == "EN" else "Doar originale",
         "88.9", "86.0", "95.9", "0.030", "-2.7"],
    ]
    abl_widths = [Inches(3.2), Inches(1.4), Inches(1.4), Inches(1.4),
                  Inches(1.4), Inches(1.2)]
    add_table_slide(slide, t["s22_title"], abl_headers, abl_rows, 25, total_slides,
                    col_widths=abl_widths)

    # -----------------------------------------------------------------------
    # Slide 26: Calibration results table
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    if lang == "RO":
        cal_headers = ["Model", "Necalibrat ECE", "Temperatura ECE",
                       "Platt ECE", "Isotonic ECE"]
    else:
        cal_headers = ["Model", "Uncalibrated ECE", "Temperature ECE",
                       "Platt ECE", "Isotonic ECE"]
    cal_rows = [
        ["XGB (ResNet-50)", "2.2%", "1.6%", "4.4%", "0.6%"],
        ["XGB (EffNet-B0)", "2.4%", "1.8%", "4.1%", "0.7%"],
        ["XGB (ViT-B/16)", "2.7%", "2.0%", "3.9%", "0.8%"],
        ["XGB (3-backbone)", "2.1%", "1.5%", "3.8%", "0.6%"],
        ["MLP (3-backbone)", "3.9%", "2.8%", "5.1%", "0.8%"],
    ]
    cal_widths = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    add_table_slide(slide, t["s23_title"], cal_headers, cal_rows, 26, total_slides,
                    col_widths=cal_widths)

    # -----------------------------------------------------------------------
    # Slide 27: Reliability diagrams figure
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    fig_path = FIGURES_DIR / "fig2_reliability_diagrams.png"
    add_figure_slide(slide, t["s24_title"], fig_path, 27, total_slides,
                     caption="Reliability diagrams before and after calibration")

    # -----------------------------------------------------------------------
    # Slide 28: Conformal prediction results
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s25_title"], t["s25_bullets"], 28, total_slides)

    # -----------------------------------------------------------------------
    # Slide 29: Statistical significance
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s26_title"], t["s26_bullets"], 29, total_slides)

    # -----------------------------------------------------------------------
    # Slide 30: Efficiency metrics
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s27_title"], t["s27_bullets"], 30, total_slides)

    # -----------------------------------------------------------------------
    # Slide 31: Overfitting diagnostic figure
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    fig_path = FIGURES_DIR / "fig9_training_curves.png"
    add_figure_slide(slide, t["s28_title"], fig_path, 31, total_slides,
                     caption="Training vs validation loss/AUC across epochs")

    # -----------------------------------------------------------------------
    # Slide 32: Engineering section divider
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    section_sub = "Partea 5: Inginerie Software" if lang == "RO" \
        else "Part 5: Software Engineering"
    add_section_slide(slide, section_sub,
                      "~100 files | 15 modules | CI/CD | Docker",
                      32, total_slides)

    # -----------------------------------------------------------------------
    # Slide 33: Codebase stats
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s29_title"], t["s29_bullets"], 33, total_slides)

    # -----------------------------------------------------------------------
    # Slide 34: Technology stack (two columns)
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_two_column_slide(
        slide, t["s30_title"],
        t["s30_bullets_left"], t["s30_bullets_right"],
        34, total_slides
    )

    # -----------------------------------------------------------------------
    # Slide 35: ML Pipeline
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s31_title"], t["s31_bullets"], 35, total_slides)

    # -----------------------------------------------------------------------
    # Slide 36: Publication artifacts
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s32_title"], t["s32_bullets"], 36, total_slides)

    # -----------------------------------------------------------------------
    # Slide 37: Future work section divider
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    section_sub = "Partea 6: Lucrari Viitoare" if lang == "RO" \
        else "Part 6: Future Work & Limitations"
    add_section_slide(slide, section_sub,
                      "Threat model | Honest limitations | Roadmap",
                      37, total_slides)

    # -----------------------------------------------------------------------
    # Slides 38-42: Future work content
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s33_title"], t["s33_bullets"], 38, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s34_title"], t["s34_bullets"], 39, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s35_title"], t["s35_bullets"], 40, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s36_title"], t["s36_bullets"], 41, total_slides)

    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide(slide, t["s37_title"], t["s37_bullets"], 42, total_slides)

    # -----------------------------------------------------------------------
    # Slide 43: Conclusions + Questions
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_text_box(
        slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
        t["s38_title"], font_size=36, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )

    # Conclusion bullets
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10), Inches(4.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(t["s38_bullets"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)

        if bullet == "":
            continue

        if bullet.startswith("  - "):
            text = bullet[4:]
            fs = Pt(16)
        elif bullet in ("Questions?", "Intrebari?"):
            text = bullet
            fs = Pt(32)
            p.space_before = Pt(20)
        else:
            text = bullet
            fs = Pt(18)

        run = p.add_run()
        run.text = text
        run.font.size = fs
        run.font.color.rgb = WHITE if text not in ("Questions?", "Intrebari?") \
            else LIGHT_BLUE
        run.font.bold = text in ("Questions?", "Intrebari?")
        run.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER if text in ("Questions?", "Intrebari?") \
            else PP_ALIGN.LEFT

    return prs


# ===========================================================================
# Main
# ===========================================================================

def main():
    print("Generating ImageTrust academic presentations...")
    print(f"Figures directory: {FIGURES_DIR}")
    print(f"Output directory: {OUTPUT_DIR}")

    # Check figures
    figs = list(FIGURES_DIR.glob("*.png"))
    print(f"Found {len(figs)} figures in {FIGURES_DIR}")

    # Update total_slides constant used in footer
    # (we actually have 43 slides due to section dividers, update footer total)

    for lang in ["RO", "EN"]:
        print(f"\nBuilding {lang} presentation...")
        prs = build_presentation(lang)

        # Update all footers to reflect actual slide count
        actual_total = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            # Fix footer numbers
                            for old_total in range(30, 50):
                                old_str = f" / {old_total}"
                                if old_str in run.text:
                                    new_num = i + 1
                                    run.text = f"{new_num} / {actual_total}"

        suffix = "RO" if lang == "RO" else "EN"
        output_path = OUTPUT_DIR / f"ImageTrust_Presentation_{suffix}.pptx"
        prs.save(str(output_path))
        print(f"Saved: {output_path} ({actual_total} slides)")

    print("\nDone! Both presentations generated successfully.")


if __name__ == "__main__":
    main()
